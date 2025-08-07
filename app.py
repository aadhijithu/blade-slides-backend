from flask import Flask, request, send_file, jsonify
from flask_cors import CORS
import io
import base64
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import os

app = Flask(__name__)
CORS(app)  # Enable CORS for Figma plugin

# Import our processing modules
from figma_processor import process_figma_data
from pptx_generator import create_presentation

@app.route('/', methods=['GET'])
def health_check():
    return jsonify({
        "status": "healthy",
        "message": "Blade Slides Backend is running",
        "version": "1.0.0"
    })

@app.route('/convert-figma-to-pptx', methods=['POST'])
def convert_figma_to_pptx():
    try:
        # Get JSON data from Figma plugin
        figma_data = request.get_json()
        
        if not figma_data:
            return jsonify({"error": "No data provided"}), 400
        
        if not figma_data.get('slides') or len(figma_data['slides']) == 0:
            return jsonify({"error": "No slides provided"}), 400
        
        # Process the Figma data and create PPTX
        print(f"Processing {len(figma_data['slides'])} slides...")
        
        # Create the presentation
        pptx_buffer = create_presentation(figma_data)
        
        # Prepare filename
        filename = figma_data.get('fileName', 'Untitled_Presentation')
        clean_filename = ''.join(c for c in filename if c.isalnum() or c in (' ', '-', '_')).rstrip()
        
        print(f"Generated PPTX: {clean_filename}")
        
        return send_file(
            pptx_buffer,
            as_attachment=True,
            download_name=f"{clean_filename}_PythonPPTX.pptx",
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
        
    except Exception as e:
        print(f"Error converting Figma to PPTX: {str(e)}")
        return jsonify({
            "error": "Failed to generate PPTX",
            "details": str(e)
        }), 500

@app.route('/test-pptx', methods=['GET'])
def test_pptx():
    """Test endpoint to verify python-pptx is working"""
    try:
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Add test text
        left = Inches(1)
        top = Inches(1)
        width = Inches(8)
        height = Inches(1)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = "Test from Railway Backend!"
        
        # Save to buffer
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        
        return send_file(
            buffer,
            as_attachment=True,
            download_name="test.pptx",
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
        
    except Exception as e:
        return jsonify({
            "error": "Test failed",
            "details": str(e)
        }), 500

if __name__ == '__main__':
    port = int(os.environ.get('PORT', 8000))
    app.run(host='0.0.0.0', port=port, debug=True) 
"""
PPTX Generator using python-pptx
High-fidelity conversion from Figma to PowerPoint with precise positioning
"""

import io
import base64
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
from figma_processor import process_figma_data
from utils.fonts import map_figma_font, map_text_alignment, map_vertical_alignment
from utils.colors import hex_to_rgb
from utils.layouts import calculate_positioning, calculate_layer_position, scale_font_size

def create_presentation(figma_data):
    """
    Create a PowerPoint presentation from Figma data
    
    Args:
        figma_data: Processed Figma data
        
    Returns:
        BytesIO buffer containing the PPTX file
    """
    # Process the Figma data
    processed_data = process_figma_data(figma_data)
    
    # Create new presentation
    prs = Presentation()
    
    # Remove default slide
    if len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]
    
    config = processed_data['config']
    
    print(f"Creating presentation with {len(processed_data['slides'])} slides...")
    
    # Process each slide
    for slide_index, slide_data in enumerate(processed_data['slides']):
        print(f"Processing slide {slide_index + 1}: {slide_data['name']}")
        
        # Add blank slide
        blank_slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Calculate precise positioning
        positioning = calculate_positioning(slide_data, config)
        
        # Set slide background
        if slide_data.get('background'):
            set_slide_background(slide, slide_data['background'], config)
        
        # Process layers (correct order for proper z-stacking)
        layers = slide_data['layers']
        sorted_layers = sorted(layers, key=lambda x: x.get('zIndex', 0), reverse=False)
        
        print(f"Processing {len(sorted_layers)} layers in z-order...")
        
        for layer in sorted_layers:
            try:
                print(f"  Processing layer: {layer['name']} (z:{layer.get('zIndex', 0)}, type:{layer['type']})")
                process_layer(slide, layer, positioning, config)
            except Exception as e:
                print(f"Error processing layer {layer['name']}: {str(e)}")
                continue
        
        # Add slide metadata
        # add_slide_metadata(slide, slide_data, slide_index + 1, len(processed_data['slides']), config)
        print(f"  Skipping slide metadata (frame name & page numbers) to match Figma design")
    
    # Save to buffer
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    
    print("PPTX generation completed successfully")
    return buffer

def process_layer(slide, layer, positioning, config):
    """
    Process a single layer and add it to the slide

    Args:
        slide: PowerPoint slide object
        layer: Layer data from Figma
        positioning: Positioning calculations
        config: Configuration settings
    """
    layer_type = layer.get('type')
    layer_name = layer.get('name', 'Unnamed Layer')
    
    # Validate layer data
    if not layer_type:
        print(f"  Warning: Skipping layer '{layer_name}' - no type specified")
        return
        
    if layer_type not in ['TEXT', 'SHAPE', 'IMAGE']:
        print(f"  Warning: Skipping layer '{layer_name}' - unknown type: {layer_type}")
        return

    if layer_type == 'TEXT':
        content = layer.get('content', '').strip()
        if not content:
            print(f"  Warning: Skipping text layer '{layer_name}' - no content")
            return
        add_text_layer(slide, layer, positioning, config)
    elif layer_type == 'SHAPE':
        add_shape_layer(slide, layer, positioning, config)
    elif layer_type == 'IMAGE':
        if not layer.get('imageData'):
            print(f"  Warning: Skipping image layer '{layer_name}' - no image data")
            return
        add_image_layer(slide, layer, positioning, config)
    
    print(f"    ✅ Successfully added {layer_type} layer: {layer_name}")

def add_text_layer(slide, layer, positioning, config):
    """
    Add a text layer to the slide with precise formatting
    """
    # Calculate position
    pos = calculate_layer_position(layer, positioning, config)
    
    # Create text box
    textbox = slide.shapes.add_textbox(
        Inches(pos['x']), 
        Inches(pos['y']), 
        Inches(pos['width']), 
        Inches(pos['height'])
    )
    
    text_frame = textbox.text_frame
    text_frame.clear()  # Remove default paragraph
    
    # Add text content
    p = text_frame.paragraphs[0]
    p.text = layer.get('content', '')
    
    # Apply text formatting
    style = layer.get('style', {})
    
    # Font and size
    font_size = style.get('fontSize', 16)
    scaled_font_size = scale_font_size(font_size, positioning['scale'])
    font_family = style.get('fontFamily', 'Arial')
    
    run = p.runs[0] if p.runs else p.runs.add()
    run.font.size = Pt(max(8, scaled_font_size))  # Minimum 8pt
    run.font.name = map_figma_font(font_family)
    
    # Font weight and style
    font_weight = style.get('fontWeight', 'Regular')
    run.font.bold = 'Bold' in font_weight or '700' in font_weight
    run.font.italic = 'Italic' in font_weight
    
    # Text color
    color_hex = style.get('color', '#000000')
    rgb = hex_to_rgb(color_hex)
    run.font.color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])
    
    # Text alignment
    text_align = style.get('textAlign', 'left')
    p.alignment = map_text_alignment(text_align)
    
    # Vertical alignment
    vertical_align = style.get('verticalAlign', 'top')
    text_frame.vertical_anchor = map_vertical_alignment(vertical_align)
    
    print(f"    Text formatting: font='{font_family}'->'{ run.font.name}', size={font_size}->{scaled_font_size:.1f}pt, color='{color_hex}', weight='{font_weight}'")
    
    # Text box properties
    text_frame.margin_left = Inches(0.05)
    text_frame.margin_right = Inches(0.05)
    text_frame.margin_top = Inches(0.05)
    text_frame.margin_bottom = Inches(0.05)
    text_frame.word_wrap = True
    text_frame.auto_size = None  # Fixed size

def add_shape_layer(slide, layer, positioning, config):
    """
    Add a shape layer to the slide
    """
    # Calculate position
    pos = calculate_layer_position(layer, positioning, config)
    
    # Get shape properties
    style = layer.get('style', {})
    shape_type = layer.get('shapeType', 'rectangle')
    
    # Create shape based on type
    if shape_type == 'rectangle':
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(pos['x']), 
            Inches(pos['y']), 
            Inches(pos['width']), 
            Inches(pos['height'])
        )
    elif shape_type == 'ellipse':
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(pos['x']), 
            Inches(pos['y']), 
            Inches(pos['width']), 
            Inches(pos['height'])
        )
    else:
        # Default to rectangle for unknown shapes
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(pos['x']), 
            Inches(pos['y']), 
            Inches(pos['width']), 
            Inches(pos['height'])
        )
    
    # Apply fill
    fill_color = style.get('fill', 'transparent')
    if fill_color and fill_color != 'transparent' and fill_color != 'none':
        rgb = hex_to_rgb(fill_color)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])
    else:
        shape.fill.background()  # Transparent fill
    
    # Apply stroke
    stroke_color = style.get('stroke')
    stroke_width = style.get('strokeWidth', 0)
    
    if stroke_color and stroke_color != 'transparent' and stroke_width > 0:
        rgb = hex_to_rgb(stroke_color)
        shape.line.color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])
        # Convert pixels to points (approximate)
        stroke_pt = max(0.25, stroke_width * positioning['scale'])
        shape.line.width = Pt(stroke_pt)
    else:
        # Explicitly remove any line/border
        shape.line.fill.background()
        
    # Remove any shadow or effects
    try:
        # Ensure no shadow effects
        if hasattr(shape, 'shadow'):
            shape.shadow.inherit = False
    except:
        pass  # Ignore if shadow property doesn't exist
    
    print(f"    Shape styling: fill='{fill_color}', stroke='{stroke_color}' ({stroke_width}px)")

def add_image_layer(slide, layer, positioning, config):
    """
    Add an image layer to the slide
    """
    image_data = layer.get('imageData')
    if not image_data:
        print(f"No image data for layer: {layer['name']}")
        return
    
    try:
        # Decode base64 image
        if image_data.startswith('data:image'):
            image_data = image_data.split(',')[1]
        
        image_bytes = base64.b64decode(image_data)
        image_stream = io.BytesIO(image_bytes)
        
        # Calculate position
        pos = calculate_layer_position(layer, positioning, config)
        
        # Add image to slide
        slide.shapes.add_picture(
            image_stream,
            Inches(pos['x']), 
            Inches(pos['y']), 
            Inches(pos['width']), 
            Inches(pos['height'])
        )
        
    except Exception as e:
        print(f"Error adding image layer {layer['name']}: {str(e)}")

def set_slide_background(slide, background, config):
    """
    Set the slide background
    """
    if not background:
        return
    
    bg_type = background.get('type')
    
    if bg_type == 'SOLID':
        color_hex = background.get('color', '#FFFFFF')
        rgb = hex_to_rgb(color_hex)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])
    
    elif bg_type == 'IMAGE' and background.get('imageData'):
        try:
            # Add background image
            image_data = background['imageData']
            if image_data.startswith('data:image'):
                image_data = image_data.split(',')[1]
            
            image_bytes = base64.b64decode(image_data)
            image_stream = io.BytesIO(image_bytes)
            
            # Add as full-slide image
            slide.shapes.add_picture(
                image_stream,
                Inches(0), 
                Inches(0), 
                Inches(config['slideWidth']), 
                Inches(config['slideHeight'])
            )
            
        except Exception as e:
            print(f"Error setting background image: {str(e)}")

def add_slide_metadata(slide, slide_data, slide_num, total_slides, config):
    """
    Add slide title and page number
    """
    # Slide title
    title_box = slide.shapes.add_textbox(
        Inches(0.2), 
        Inches(0.05), 
        Inches(9.6), 
        Inches(0.3)
    )
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = slide_data['name']
    
    title_run = title_p.runs[0]
    title_run.font.size = Pt(10)
    title_run.font.name = 'Arial'
    title_run.font.color.rgb = RGBColor(102, 102, 102)  # #666666
    title_run.font.bold = True
    
    # Page number
    page_box = slide.shapes.add_textbox(
        Inches(8.5), 
        Inches(config['slideHeight'] - 0.4), 
        Inches(1.3), 
        Inches(0.3)
    )
    page_frame = page_box.text_frame
    page_p = page_frame.paragraphs[0]
    page_p.text = f"{slide_num} / {total_slides}"
    page_p.alignment = PP_ALIGN.RIGHT
    
    page_run = page_p.runs[0]
    page_run.font.size = Pt(8)
    page_run.font.name = 'Arial'
    page_run.font.color.rgb = RGBColor(136, 136, 136)  # #888888 
